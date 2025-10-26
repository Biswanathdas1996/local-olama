"""
Enterprise-grade document extraction using Docling for structure-aware parsing.
Supports PDF, DOCX, TXT, PPTX, HTML with metadata preservation.
"""

import io
import json
import logging
from pathlib import Path
from typing import Dict, List, Optional, BinaryIO
from dataclasses import dataclass
from datetime import datetime

try:
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    DOCLING_AVAILABLE = True
except ImportError as e:
    # Fallback to basic extraction if docling is not available
    print(f"⚠️  Docling import failed: {e}")
    import traceback
    traceback.print_exc()
    DocumentConverter = None
    PdfFormatOption = None
    InputFormat = None
    PdfPipelineOptions = None
    DOCLING_AVAILABLE = False
except Exception as e:
    # Catch any other exceptions during import
    print(f"❌ Unexpected error importing Docling: {e}")
    import traceback
    traceback.print_exc()
    DocumentConverter = None
    PdfFormatOption = None
    InputFormat = None
    PdfPipelineOptions = None
    DOCLING_AVAILABLE = False

from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


@dataclass
class ExtractedDocument:
    """Structured document extraction result"""
    text: str
    metadata: Dict[str, any]
    sections: List[Dict[str, str]]  # [{title, content, page}]
    format: str
    images: List[Dict[str, any]] = None  # List of extracted images with metadata
    
    def __post_init__(self):
        """Initialize images list if None"""
        if self.images is None:
            self.images = []


class DocumentExtractor:
    """
    High-fidelity document extractor using Docling for structure-aware parsing.
    Falls back to format-specific parsers when needed.
    """

    SUPPORTED_FORMATS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.txt': 'text',
        '.pptx': 'pptx',
        '.html': 'html',
        '.htm': 'html',
    }

    def __init__(self, use_ocr: bool = False, output_dir: Optional[str] = None):
        """
        Initialize document extractor.
        
        Args:
            use_ocr: Enable OCR for scanned PDFs (requires tesseract)
            output_dir: Directory to save Docling extraction outputs (default: data/docling_output)
        """
        self.use_ocr = use_ocr
        self.converter = None
        
        # Setup output directory for Docling results
        if output_dir:
            self.output_dir = Path(output_dir)
        else:
            self.output_dir = Path("data/docling_output")
        
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Docling output directory: {self.output_dir}")
        logger.info(f"Docling availability check: DOCLING_AVAILABLE={DOCLING_AVAILABLE}, DocumentConverter={DocumentConverter is not None}")
        
        if DOCLING_AVAILABLE:
            try:
                # Configure Docling with the new API (v2.x)
                # Create format options for PDF with OCR settings
                format_options = {}
                
                if PdfFormatOption and PdfPipelineOptions:
                    pipeline_options = PdfPipelineOptions()
                    pipeline_options.do_ocr = use_ocr
                    format_options[InputFormat.PDF] = PdfFormatOption(
                        pipeline_options=pipeline_options
                    )
                
                self.converter = DocumentConverter(
                    allowed_formats=[
                        InputFormat.PDF,
                        InputFormat.DOCX,
                        InputFormat.HTML,
                        InputFormat.PPTX,
                    ],
                    format_options=format_options if format_options else None,
                )
                logger.info("✅ Docling converter initialized successfully")
                print("\n" + "="*80)
                print("✅ DOCLING INITIALIZED")
                print(f"   OCR Enabled: {use_ocr}")
                print(f"   Output Directory: {self.output_dir}")
                print(f"   Supported Formats: PDF, DOCX, HTML, PPTX")
                print("="*80 + "\n")
            except Exception as e:
                logger.warning(f"❌ Docling initialization failed: {e}. Using fallback extractors.")
                print("\n" + "!"*80)
                print("⚠️  DOCLING INITIALIZATION FAILED")
                print(f"   Error: {str(e)[:100]}")
                print("   Falling back to basic extractors")
                print("!"*80 + "\n")
                self.converter = None
        else:
            logger.warning("❌ Docling not available. Using fallback extractors.")
            print("\n" + "!"*80)
            print("⚠️  DOCLING NOT INSTALLED")
            print("   Install with: pip install docling")
            print("   Falling back to basic extractors")
            print("!"*80 + "\n")

    def extract(
        self, 
        file_content: BinaryIO, 
        filename: str,
        **kwargs
    ) -> ExtractedDocument:
        """
        Extract structured content from document.
        
        Args:
            file_content: Binary file content
            filename: Original filename
            **kwargs: Additional extraction parameters
            
        Returns:
            ExtractedDocument with text, metadata, and sections
        """
        file_ext = Path(filename).suffix.lower()
        
        if file_ext not in self.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {file_ext}. Supported: {list(self.SUPPORTED_FORMATS.keys())}")
        
        doc_format = self.SUPPORTED_FORMATS[file_ext]
        
        # Try Docling first for supported formats
        if self.converter and doc_format in ['pdf', 'docx', 'html', 'pptx']:
            try:
                logger.warning(f"⚡ DOCLING EXTRACTION: Processing '{filename}' with Docling for high-fidelity extraction...")
                print(f"\n{'='*80}")
                print(f"⚡ DOCLING EXTRACTION ACTIVE")
                print(f"   File: {filename}")
                print(f"   Format: {doc_format.upper()}")
                print(f"   Output will be saved to: {self.output_dir}")
                print(f"{'='*80}\n")
                return self._extract_with_docling(file_content, filename, doc_format)
            except Exception as e:
                logger.warning(f"Docling extraction failed for {filename}: {e}. Trying fallback.")
                print(f"\n{'!'*80}")
                print(f"⚠️  DOCLING EXTRACTION FAILED - USING FALLBACK")
                print(f"   File: {filename}")
                print(f"   Error: {str(e)[:100]}")
                print(f"   Switching to basic {doc_format.upper()} extractor...")
                print(f"{'!'*80}\n")
        else:
            # Log when Docling is not available or format not supported
            if not self.converter:
                logger.warning(f"⚠️  FALLBACK EXTRACTION: Docling not available. Using basic extractor for '{filename}'")
                print(f"\n{'!'*80}")
                print(f"⚠️  FALLBACK EXTRACTION (Docling Not Available)")
                print(f"   File: {filename}")
                print(f"   Format: {doc_format.upper()}")
                print(f"   Using basic extractor - limited structure preservation")
                print(f"{'!'*80}\n")
            else:
                logger.warning(f"⚠️  FALLBACK EXTRACTION: Format '{doc_format}' not supported by Docling. Using basic extractor for '{filename}'")
                print(f"\n{'!'*80}")
                print(f"⚠️  FALLBACK EXTRACTION (Format Not Supported)")
                print(f"   File: {filename}")
                print(f"   Format: {doc_format.upper()}")
                print(f"   Using basic extractor - limited structure preservation")
                print(f"{'!'*80}\n")
        
        # Fallback to format-specific extractors
        extractors = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'text': self._extract_text,
            'pptx': self._extract_pptx,
            'html': self._extract_html,
        }
        
        extractor = extractors.get(doc_format)
        if not extractor:
            raise ValueError(f"No extractor available for format: {doc_format}")
        
        return extractor(file_content, filename)

    def _extract_with_docling(
        self, 
        file_content: BinaryIO, 
        filename: str,
        doc_format: str
    ) -> ExtractedDocument:
        """Extract using Docling for maximum structure preservation"""
        
        # Save to temporary file for Docling
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as tmp:
            tmp.write(file_content.read())
            tmp_path = tmp.name
        
        try:
            # Convert document
            result = self.converter.convert(tmp_path)
            
            # Create output subdirectory for this document
            doc_name = Path(filename).stem
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            doc_output_dir = self.output_dir / f"{doc_name}_{timestamp}"
            doc_output_dir.mkdir(parents=True, exist_ok=True)
            
            # Create images subdirectory
            images_dir = doc_output_dir / "images"
            images_dir.mkdir(exist_ok=True)
            
            # Extract images if available
            extracted_images = []
            try:
                # First, try to get pictures from the document's dict export
                # (Sometimes result.document.pictures is empty but the JSON has them)
                doc_dict = result.document.export_to_dict()
                
                if 'pictures' in doc_dict and doc_dict['pictures']:
                    logger.info(f"✅ Found {len(doc_dict['pictures'])} images in document JSON export")
                    
                    for idx, picture in enumerate(doc_dict['pictures']):
                        try:
                            # Save image
                            image_id = f"{doc_name}_img_{idx}"
                            image_path = images_dir / f"{image_id}.png"
                            
                            # Get image data from JSON structure
                            if 'image' in picture and picture['image'] and 'uri' in picture['image']:
                                import base64
                                import io
                                from PIL import Image
                                
                                uri = picture['image']['uri']
                                if uri.startswith('data:image'):
                                    # Extract base64 data from URI
                                    base64_data = uri.split(',')[1] if ',' in uri else uri
                                    img_bytes = base64.b64decode(base64_data)
                                    img = Image.open(io.BytesIO(img_bytes))
                                    img.save(image_path, 'PNG')
                                    
                                    # Extract metadata
                                    image_metadata = {
                                        'image_id': image_id,
                                        'image_index': idx,
                                        'source_document': filename,
                                        'image_path': str(image_path),
                                        'mimetype': picture['image'].get('mimetype'),
                                        'dpi': picture['image'].get('dpi'),
                                    }
                                    
                                    # Add captions if available
                                    if 'captions' in picture and picture['captions']:
                                        image_metadata['caption'] = ' '.join([c.get('text', '') for c in picture['captions']])
                                    
                                    extracted_images.append(image_metadata)
                                    logger.info(f"✅ Saved image {idx} to {image_path}")
                                else:
                                    logger.warning(f"Image {idx} has non-data URI: {uri[:50]}")
                            else:
                                logger.warning(f"Image {idx} has no accessible URI data")
                                
                        except Exception as img_err:
                            logger.warning(f"Failed to extract image {idx}: {img_err}")
                
                # Fallback: Try the old method with result.document.pictures
                elif hasattr(result.document, 'pictures') and result.document.pictures:
                    logger.info(f"Found {len(result.document.pictures)} images in document.pictures")
                    for idx, picture in enumerate(result.document.pictures):
                        try:
                            # Save image
                            image_id = f"{doc_name}_img_{idx}"
                            image_path = images_dir / f"{image_id}.png"
                            
                            # Get image data - different APIs might use different attributes
                            import base64
                            import io
                            from PIL import Image
                            
                            saved = False
                            
                            # Try different image attribute patterns
                            if hasattr(picture, 'image') and picture.image is not None:
                                # Check if it's a URI (base64 data)
                                if hasattr(picture.image, 'uri') and picture.image.uri:
                                    uri = picture.image.uri
                                    if uri.startswith('data:image'):
                                        # Extract base64 data from URI
                                        base64_data = uri.split(',')[1] if ',' in uri else uri
                                        img_bytes = base64.b64decode(base64_data)
                                        img = Image.open(io.BytesIO(img_bytes))
                                        img.save(image_path, 'PNG')
                                        saved = True
                                # PIL image object
                                elif hasattr(picture.image, 'save'):
                                    picture.image.save(image_path, 'PNG')
                                    saved = True
                            elif hasattr(picture, 'pil_image') and picture.pil_image is not None:
                                picture.pil_image.save(image_path, 'PNG')
                                saved = True
                            elif hasattr(picture, 'data'):
                                # Save raw data
                                image_path.write_bytes(picture.data)
                                saved = True
                            
                            if not saved:
                                logger.warning(f"Image {idx} has no accessible data")
                                continue
                            
                            # Extract metadata
                            image_metadata = {
                                'image_id': image_id,
                                'image_index': idx,
                                'source_document': filename,
                                'image_path': str(image_path),
                                'page': getattr(picture, 'page', None),
                                'caption': getattr(picture, 'caption', None) or getattr(picture, 'text', None),
                            }
                            
                            # Add bounding box if available
                            if hasattr(picture, 'bbox'):
                                image_metadata['bbox'] = picture.bbox
                            
                            extracted_images.append(image_metadata)
                            logger.info(f"Saved image {idx} to {image_path}")
                            
                        except Exception as img_err:
                            logger.warning(f"Failed to extract image {idx}: {img_err}")
                
                else:
                    logger.info("ℹ️  No images found in document")
                
            except Exception as img_extract_err:
                logger.warning(f"Image extraction failed: {img_extract_err}")
            
            # Save Docling outputs in multiple formats
            try:
                # 1. Save as Markdown
                markdown_content = result.document.export_to_markdown()
                markdown_path = doc_output_dir / f"{doc_name}.md"
                markdown_path.write_text(markdown_content, encoding='utf-8')
                logger.info(f"Saved Markdown: {markdown_path}")
                
                # 2. Save as JSON (structured data)
                try:
                    json_content = result.document.export_to_dict()
                    json_path = doc_output_dir / f"{doc_name}.json"
                    json_path.write_text(json.dumps(json_content, indent=2, ensure_ascii=False, default=str), encoding='utf-8')
                    logger.info(f"Saved JSON: {json_path}")
                except Exception as json_err:
                    logger.warning(f"Could not save JSON export: {json_err}")
                
                # 3. Save as plain text
                text_content = result.document.export_to_text()
                text_path = doc_output_dir / f"{doc_name}.txt"
                text_path.write_text(text_content, encoding='utf-8')
                logger.info(f"Saved Text: {text_path}")
                
                # 4. Save metadata (simplified to avoid serialization issues)
                metadata_info = {
                    'filename': filename,
                    'format': doc_format,
                    'extractor': 'docling',
                    'extraction_timestamp': timestamp,
                    'output_directory': str(doc_output_dir),
                    'images_extracted': len(extracted_images),
                    'exports': {
                        'markdown': str(markdown_path),
                        'json': str(json_path),
                        'text': str(text_path)
                    }
                }
                metadata_path = doc_output_dir / f"{doc_name}_metadata.json"
                metadata_path.write_text(json.dumps(metadata_info, indent=2), encoding='utf-8')
                logger.info(f"Saved Metadata: {metadata_path}")
                
                # Success message
                print(f"\n{'='*80}")
                print(f"✅ DOCLING EXTRACTION SUCCESSFUL")
                print(f"   File: {filename}")
                print(f"   Pages: {getattr(result.document, 'num_pages', 'N/A')}")
                print(f"   Images: {len(extracted_images)}")
                print(f"   Output Directory: {doc_output_dir}")
                print(f"   Saved Formats:")
                print(f"      ├── Markdown: {doc_name}.md")
                print(f"      ├── JSON: {doc_name}.json")
                print(f"      ├── Text: {doc_name}.txt")
                print(f"      ├── Metadata: {doc_name}_metadata.json")
                if extracted_images:
                    print(f"      └── Images: {len(extracted_images)} saved in images/")
                print(f"{'='*80}\n")
                
            except Exception as e:
                logger.error(f"Error saving Docling outputs: {e}")
                print(f"\n{'!'*80}")
                print(f"⚠️  WARNING: Docling extracted content but failed to save outputs")
                print(f"   Error: {str(e)[:100]}")
                print(f"   Extraction will continue with in-memory content")
                print(f"{'!'*80}\n")
            
            # Extract structured content
            text = markdown_content
            
            # Replace image placeholders with OCR-extracted content if images were found
            if extracted_images:
                text = self._inject_image_descriptions(text, extracted_images)
            
            # Extract sections with hierarchy (with error handling for API changes)
            sections = []
            try:
                # Try to iterate through document elements
                if hasattr(result.document, 'body') and hasattr(result.document.body, 'children'):
                    # Docling 2.x API
                    logger.info(f"Extracting sections from {len(result.document.body.children)} body children")
                    for element in result.document.body.children:
                        if hasattr(element, 'label') and element.label and 'heading' in str(element.label).lower():
                            sections.append({
                                'title': getattr(element, 'text', str(element))[:100],
                                'content': '',
                                'level': getattr(element, 'level', 1)
                            })
                        elif sections and hasattr(element, 'text'):
                            sections[-1]['content'] += getattr(element, 'text', '') + '\n'
                elif hasattr(result.document, 'body') and hasattr(result.document.body, 'elements'):
                    # Older API
                    logger.info(f"Extracting sections from {len(result.document.body.elements)} body elements")
                    for element in result.document.body.elements:
                        if hasattr(element, 'heading_level') and element.heading_level:
                            sections.append({
                                'title': element.text,
                                'content': '',
                                'level': element.heading_level,
                                'page': getattr(element, 'page', None)
                            })
                        elif sections and hasattr(element, 'text'):
                            sections[-1]['content'] += element.text + '\n'
                
                # If no sections found, create sections from pages/slides
                if not sections and doc_format == 'pptx':
                    logger.info("No sections found via standard extraction, trying page-based extraction for PPTX")
                    # For PPTX, each page might be a slide
                    if hasattr(result.document, 'pages'):
                        for page_num, page in enumerate(result.document.pages, 1):
                            page_text = ""
                            if hasattr(page, 'text'):
                                page_text = page.text
                            elif hasattr(page, 'export_to_text'):
                                page_text = page.export_to_text()
                            
                            if page_text.strip():
                                sections.append({
                                    'title': f'Slide {page_num}',
                                    'content': page_text,
                                    'slide': page_num
                                })
                        logger.info(f"Created {len(sections)} sections from {len(result.document.pages)} pages")
                    
            except Exception as section_err:
                logger.warning(f"Could not extract sections: {section_err}")
                import traceback
                traceback.print_exc()
                # Create a simple section from the whole text
                sections = [{'title': filename, 'content': text}]
            
            # Final fallback: if still no sections, create one from text
            if not sections and text.strip():
                logger.warning(f"No sections extracted, creating default section from text content")
                sections = [{'title': filename, 'content': text}]
            
            metadata = {
                'filename': filename,
                'format': doc_format,
                'extractor': 'docling',
                'output_directory': str(doc_output_dir),
                'images_count': len(extracted_images),
            }
            
            return ExtractedDocument(
                text=text,
                metadata=metadata,
                sections=sections,
                format=doc_format,
                images=extracted_images
            )
            
        finally:
            # Cleanup
            Path(tmp_path).unlink(missing_ok=True)
    
    def _inject_image_descriptions(self, text: str, extracted_images: List[Dict]) -> str:
        """
        Replace <!-- image --> placeholders with OCR-extracted image content.
        Uses quick OCR to get text from images and injects it inline.
        """
        logger.info(f"🖼️  _inject_image_descriptions called: {len(extracted_images) if extracted_images else 0} images")
        logger.info(f"🔍 Text contains <!-- image -->: {'<!-- image -->' in text}")
        
        if not extracted_images or "<!-- image -->" not in text:
            logger.warning("⚠️  Skipping image injection - no images or no placeholders found")
            return text
        
        try:
            from core.image_processor import get_image_processor
            img_processor = get_image_processor()
            
            logger.info(f"Processing {len(extracted_images)} images for inline descriptions...")
            
            # Process each image
            for idx, image_meta in enumerate(extracted_images):
                try:
                    # Quick OCR processing (no chart parsing for speed)
                    image_content = img_processor.process_image(
                        image_source=image_meta['image_path'],
                        image_id=image_meta['image_id'],
                        metadata=image_meta,
                        save_output=False  # Don't save duplicates
                    )
                    
                    # Create replacement text with image description
                    replacement = f"\n\n[IMAGE: {image_meta.get('image_id', f'img_{idx}')}]\n"
                    
                    # Add OCR text if available
                    if image_content.ocr_text:
                        replacement += f"Text in image: {image_content.ocr_text}\n"
                    
                    # Add chart/table data if available
                    if image_content.structured_data:
                        if image_content.image_type == 'chart':
                            chart_desc = self._describe_chart_inline(image_content.structured_data)
                            if chart_desc:
                                replacement += f"Chart content: {chart_desc}\n"
                        elif image_content.image_type == 'table':
                            table_desc = self._describe_table_inline(image_content.structured_data)
                            if table_desc:
                                replacement += f"Table content: {table_desc}\n"
                    
                    replacement += f"[END IMAGE]\n\n"
                    
                    # Replace first occurrence of <!-- image -->
                    text = text.replace("<!-- image -->", replacement, 1)
                    logger.info(f"Injected description for image {image_meta.get('image_id')}")
                    
                except Exception as img_err:
                    logger.warning(f"Failed to process image {image_meta.get('image_id')}: {img_err}")
                    # Replace with simple placeholder
                    text = text.replace("<!-- image -->", f"\n[Image: {image_meta.get('image_id', 'unknown')}]\n", 1)
            
            return text
            
        except Exception as e:
            logger.error(f"Image description injection failed: {e}")
            return text  # Return original text if processing fails
    
    def _describe_chart_inline(self, chart_data: Dict) -> str:
        """Generate brief inline description of chart"""
        parts = []
        if chart_data.get('title'):
            parts.append(chart_data['title'])
        if chart_data.get('text_elements'):
            parts.append(' '.join(chart_data['text_elements'][:10]))  # First 10 elements
        return ' - '.join(parts) if parts else ''
    
    def _describe_table_inline(self, table_data: Dict) -> str:
        """Generate brief inline description of table"""
        if table_data.get('cells'):
            return ' | '.join(table_data['cells'][:20])  # First 20 cells
        return ''

    def _extract_pdf(self, file_content: BinaryIO, filename: str) -> ExtractedDocument:
        """Fallback PDF extraction using PyPDF2"""
        reader = PdfReader(file_content)
        
        text = ""
        sections = []
        
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            text += page_text + "\n\n"
            
            sections.append({
                'title': f'Page {page_num}',
                'content': page_text,
                'page': page_num
            })
        
        metadata = {
            'filename': filename,
            'format': 'pdf',
            'num_pages': len(reader.pages),
            'extractor': 'pypdf2',
        }
        
        return ExtractedDocument(
            text=text, 
            metadata=metadata, 
            sections=sections, 
            format='pdf',
            images=[]  # No image extraction in fallback
        )

    def _extract_docx(self, file_content: BinaryIO, filename: str) -> ExtractedDocument:
        """Extract from DOCX with paragraph structure and images"""
        doc = DocxDocument(file_content)
        
        text = ""
        sections = []
        current_section = None
        extracted_images = []
        
        # Create output directory for images
        doc_name = Path(filename).stem
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        doc_output_dir = self.output_dir / f"{doc_name}_{timestamp}"
        images_dir = doc_output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)
        
        image_counter = 0
        
        # Extract text
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if not para_text:
                continue
            
            # Detect headings
            if para.style.name.startswith('Heading'):
                if current_section:
                    sections.append(current_section)
                current_section = {
                    'title': para_text,
                    'content': '',
                    'level': int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                }
            elif current_section:
                current_section['content'] += para_text + '\n'
            
            text += para_text + '\n'
        
        if current_section:
            sections.append(current_section)
        
        # Extract images from relationships
        try:
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                        image_id = f"{doc_name}_img_{image_counter}"
                        
                        # Determine extension from content type
                        ext_map = {
                            'image/png': 'png',
                            'image/jpeg': 'jpg',
                            'image/jpg': 'jpg',
                            'image/gif': 'gif',
                            'image/bmp': 'bmp',
                            'image/tiff': 'tiff',
                        }
                        ext = ext_map.get(image_part.content_type, 'png')
                        
                        image_path = images_dir / f"{image_id}.{ext}"
                        
                        # Save image
                        with open(image_path, 'wb') as f:
                            f.write(image_part.blob)
                        
                        # Extract metadata
                        image_metadata = {
                            'image_id': image_id,
                            'image_index': image_counter,
                            'source_document': filename,
                            'image_path': str(image_path),
                            'image_format': ext,
                            'content_type': image_part.content_type,
                            'caption': None,
                        }
                        
                        extracted_images.append(image_metadata)
                        image_counter += 1
                        logger.info(f"Extracted image {image_counter} from DOCX")
                        
                    except Exception as img_err:
                        logger.warning(f"Failed to extract image: {img_err}")
            
            if extracted_images:
                logger.info(f"Extracted {len(extracted_images)} images from {filename}")
                print(f"\n{'='*80}")
                print(f"✅ DOCX IMAGE EXTRACTION")
                print(f"   File: {filename}")
                print(f"   Images: {len(extracted_images)}")
                print(f"   Output: {images_dir}")
                print(f"{'='*80}\n")
                
        except Exception as e:
            logger.warning(f"Image extraction from DOCX failed: {e}")
        
        metadata = {
            'filename': filename,
            'format': 'docx',
            'extractor': 'python-docx',
            'images_count': len(extracted_images),
            'output_directory': str(doc_output_dir) if extracted_images else None,
        }
        
        return ExtractedDocument(
            text=text, 
            metadata=metadata, 
            sections=sections, 
            format='docx',
            images=extracted_images
        )

    def _extract_text(self, file_content: BinaryIO, filename: str) -> ExtractedDocument:
        """Extract plain text"""
        text = file_content.read().decode('utf-8', errors='ignore')
        
        metadata = {
            'filename': filename,
            'format': 'text',
            'extractor': 'plain',
        }
        
        sections = [{'title': filename, 'content': text}]
        
        return ExtractedDocument(
            text=text, 
            metadata=metadata, 
            sections=sections, 
            format='text',
            images=[]
        )

    def _extract_pptx(self, file_content: BinaryIO, filename: str) -> ExtractedDocument:
        """Extract from PowerPoint with slide structure and images"""
        prs = Presentation(file_content)
        
        text = ""
        sections = []
        extracted_images = []
        
        # Create output directory for images
        doc_name = Path(filename).stem
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        doc_output_dir = self.output_dir / f"{doc_name}_{timestamp}"
        images_dir = doc_output_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)
        
        image_counter = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            slide_title = f"Slide {slide_num}"
            
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()
                    if shape_text:
                        slide_text += shape_text + '\n'
                        # First text is usually title
                        if not slide_title.startswith("Slide"):
                            slide_title = shape_text[:50]
                
                # Extract images
                try:
                    if hasattr(shape, "image"):
                        # This shape contains an image
                        image = shape.image
                        image_id = f"{doc_name}_slide{slide_num}_img_{image_counter}"
                        image_path = images_dir / f"{image_id}.{image.ext}"
                        
                        # Save image
                        with open(image_path, 'wb') as f:
                            f.write(image.blob)
                        
                        # Extract metadata
                        image_metadata = {
                            'image_id': image_id,
                            'image_index': image_counter,
                            'source_document': filename,
                            'image_path': str(image_path),
                            'slide': slide_num,
                            'slide_title': slide_title,
                            'image_format': image.ext,
                            'caption': None,  # PPTX doesn't have captions in shapes
                        }
                        
                        extracted_images.append(image_metadata)
                        image_counter += 1
                        logger.info(f"Extracted image {image_counter} from slide {slide_num}")
                        
                except AttributeError:
                    # Shape doesn't have an image
                    pass
                except Exception as img_err:
                    logger.warning(f"Failed to extract image from slide {slide_num}: {img_err}")
            
            text += slide_text + "\n\n"
            sections.append({
                'title': slide_title,
                'content': slide_text,
                'slide': slide_num
            })
        
        if extracted_images:
            logger.info(f"Extracted {len(extracted_images)} images from {filename}")
            print(f"\n{'='*80}")
            print(f"✅ PPTX IMAGE EXTRACTION")
            print(f"   File: {filename}")
            print(f"   Slides: {len(prs.slides)}")
            print(f"   Images: {len(extracted_images)}")
            print(f"   Output: {images_dir}")
            print(f"{'='*80}\n")
        
        metadata = {
            'filename': filename,
            'format': 'pptx',
            'num_slides': len(prs.slides),
            'extractor': 'python-pptx',
            'images_count': len(extracted_images),
            'output_directory': str(doc_output_dir),
        }
        
        return ExtractedDocument(
            text=text, 
            metadata=metadata, 
            sections=sections, 
            format='pptx',
            images=extracted_images
        )

    def _extract_html(self, file_content: BinaryIO, filename: str) -> ExtractedDocument:
        """Extract from HTML with structure preservation"""
        html = file_content.read()
        soup = BeautifulSoup(html, 'lxml')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        text = soup.get_text(separator='\n', strip=True)
        
        # Extract sections from headings
        sections = []
        for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            level = int(heading.name[1])
            content = []
            
            # Get content until next heading
            for sibling in heading.find_next_siblings():
                if sibling.name and sibling.name.startswith('h') and sibling.name[1].isdigit():
                    break
                if sibling.get_text(strip=True):
                    content.append(sibling.get_text(strip=True))
            
            sections.append({
                'title': heading.get_text(strip=True),
                'content': '\n'.join(content),
                'level': level
            })
        
        metadata = {
            'filename': filename,
            'format': 'html',
            'extractor': 'beautifulsoup4',
        }
        
        return ExtractedDocument(
            text=text, 
            metadata=metadata, 
            sections=sections, 
            format='html',
            images=[]
        )
